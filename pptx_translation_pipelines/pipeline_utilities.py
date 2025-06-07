import os
import tempfile
import uuid
import zipfile
import shutil
from lxml import etree
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Tuple, Optional, Generator, List, Dict, Any, Callable, Set
import copy # Needed for deep copying elements
import re # For parsing adjustment values
from pathlib import Path # For easier path manipulation
import requests
import json
import firebase_admin
from datetime import timedelta
import traceback
import zipfile
import math
from lxml import etree
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE
from firebase_admin import credentials, storage, db
import boto3
import tempfile

db_url = 'https://snb-ai-translation-agent-default-rtdb.firebaseio.com'
secret = 'nAWmdbcHRL9UGDOP0S1Rp0pZ2c7TEIapUrsEBzHJ'
download_folder = "downloads"
# Path to your service account key

def get_service_account_key():
    """Retrieve service account key from AWS Secrets Manager"""
    secret_name = "firebase/adminSdk/privateKeys"
    region_name = os.environ.get('AWS_REGION', 'us-east-1')
    
    # Create a Secrets Manager client
    session = boto3.session.Session()
    client = session.client(
        service_name='secretsmanager',
        region_name=region_name
    )
    
    try:
        get_secret_value_response = client.get_secret_value(SecretId=secret_name)
        secret = get_secret_value_response['SecretString']
        return json.loads(secret)['aita-pipeline-ec2-1']
    except Exception as e:
        print(f"Error retrieving secret: {e}")
        raise

# Get service account credentials from AWS Secrets Manager
service_account_info = get_service_account_key()

# Write the service account info to a temporary JSON file
temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False)
with temp_file as f:
    json.dump(service_account_info, f)
    temp_file_path = f.name

# Use the temporary file path with Firebase credentials
cred = credentials.Certificate(temp_file_path)

# Clean up the temporary file after initialization (optional, but good practice)
# Note: We keep it for now since Firebase might need to re-read it
# os.unlink(temp_file_path)

firebase_initialized = False

def init_firebase():
    global firebase_initialized
    if not firebase_initialized:
        try:
        # Initialize app with storage bucket
            firebase_admin.initialize_app(cred, {
                'storageBucket': 'snb-ai-translation-agent.firebasestorage.app',
                'databaseURL': db_url
            })
            firebase_initialized = True
        except Exception as e:
            pass


class Logger:
    def __init__(self, id: str = "0"):
        self.publish_ref = f"logs/{id}"

    def info(self, message):
        print(f"INFO: {message}")

    def debug(self, message):
        print(f"DEBUG: {message}")

    def warning(self, message):
        self.print_and_write(f"WARNING: {message}")

    def error(self, message):
        self.print_and_write(f"ERROR: {message}")

    def exception(self, message):
        self.print_and_write(f"EXCEPTION: {message}")

    def print_and_write(self, message):
        print(message)
        if not os.path.exists("logs"):
            os.makedirs("logs")
        with open(f"{self.publish_ref}.txt", "a") as f:
            f.write(message + "\n")

    def publish(self, message):
        init_firebase()
        try:
            ref = db.reference(f"/{self.publish_ref}")
            # Push a new object (auto-generates a unique key)
            ref.push({
                'message': str(message)
            })
        except Exception as e:
            self.exception(e)
        print(message)


def upload_output(output_path: str):
    init_firebase()
    bucket = storage.bucket()
    blob = bucket.blob(output_path)

    # Upload from local file
    blob.upload_from_filename(output_path)

    # Optionally make it publicly accessible
    return blob.generate_signed_url(expiration=timedelta(hours=2))
    blob.make_public()
    return blob.public_url


def download_file(url: str, filename: str):
    response = requests.get(url)
    if not os.path.exists(download_folder):
        os.makedirs(download_folder)
    output_path = f"{download_folder}/{filename}"
    # Save the content to a file
    with open(output_path, 'wb') as file:
        file.write(response.content)
    return output_path


def clear_id(id: str):
    requests.put(
        f"{db_url}/translation_requests/{id}.json?auth={secret}",
        data=json.dumps({})
    )
