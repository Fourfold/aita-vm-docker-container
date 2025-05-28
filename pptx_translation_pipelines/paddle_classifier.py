from paddlex import create_model
from pdf2image import convert_from_path
import traceback
import os
import zipfile
import math
from lxml import etree
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pipeline_utilities import Logger

paddle_text_types = {
    "doc_title": "Page Title",
    "paragraph_title": "Paragraph Header",
    "text": "Body",
    "page_number": "Formals",
    "abstract": "Body",
    "table_of_contents": "Page Title",
    "references": "Formals",
    "footnotes": "Formals",
    "header": "Formals",
    "footer": "Formals",
    "algorithm": "Unknown",
    "formula": "Unknown",
    "formula_number": "Formals",
    # "image": "Shape",
    "figure_caption": "Paragraph Header",
    "table": "Table",
    "table_caption": "Paragraph Header",
    "seal": "Formals",
    "figure_title": "Paragraph Header",
    # "figure": "Shape",
    "header_image": "Formals",
    "footer_image": "Formals",
    "sidebar_text": "Body"
}

class LayoutClassifier:
    _instance = None

    def __new__(cls, *args, **kwargs):
        return cls.initialize()
    

    def initialize(cls):
        if cls._instance is None:
            cls._instance = super(LayoutClassifier, cls).__new__(cls)
            cls._instance.model = create_model(model_name="PP-DocLayout-L")
        return cls._instance


    def get_source(self, ppt_path: str, pdf_path: str, id: str):
        logger = Logger(id)

        def convert_pdf_to_images_sm(pdf_path, output_folder='pdf2image_output'):
            """Converts a PDF file to image files in the SageMaker environment."""
            if not os.path.exists(output_folder):
                os.makedirs(output_folder)
            if not os.path.exists(pdf_path):
                logger.error("PDF does not exist.")
                return []

            image_paths = []
            logger.publish("Converting pdf to images...")
            try:
                # In SageMaker's Linux environment, poppler should be found if installed via apt-get
                images = convert_from_path(pdf_path)
                for i, img in enumerate(images):
                    img_path = os.path.join(output_folder, f'page_{i+1}.png')
                    img.save(img_path, 'PNG')
                    image_paths.append(img_path)
                logger.publish(f"Converted {len(image_paths)} pages to images.")
                return image_paths
            except Exception as e:
                logger.error(f"Error converting PDF: {e}")
                return []


        def analyze_document_layout(pdf_path, number_of_slides, output_save_dir='paddle_output'):
            """
            Converts PDF to images and runs layout analysis using the PPStructure engine.
            Focuses on extracting layout information (type, bbox).
            """
            image_files = convert_pdf_to_images_sm(pdf_path)
            if not image_files:
                logger.publish("PDF processing failed.")
                return None

            if not os.path.exists(output_save_dir):
                os.makedirs(output_save_dir)

            all_page_layout_results = []

            for i, img_path in enumerate(image_files):
                logger.publish(f"Processing page #{i+1} of {number_of_slides}")
                output = self.model.predict(img_path, batch_size=1, layout_nms=True)
                page_layout_info = []
                size = None
                for res in output:
                    if size is None:
                        size = res['input_img'].shape
                    page_layout_info.append(res['boxes'])
                    # # Print result and save images
                    # res.print()
                    # res.save_to_img(save_path=output_save_dir)
                    # res.save_to_json(save_path=f"{output_save_dir}/res_{img_name}.json")

                all_page_layout_results.append({
                    'image_path': img_path,
                    'image_size': size[0:2],
                    'layout_results': page_layout_info
                })
                os.remove(img_path)

            logger.publish("Document layout analysis complete.")
            return all_page_layout_results


        def apply_layout_types(extracted_shapes, layout):
            number_of_slides = len(extracted_shapes)
            source = []
            for i in range(number_of_slides):
                logger.publish(f"Processing slide #{i + 1} of {number_of_slides}")
                slide_source = []
                slide = extracted_shapes[i]
                for nshape in slide:  # [x1,x2,y1,y2]
                    size = layout[i]['image_size']  # (height,width) in pixels
                    shape = [size[1]*nshape[0], size[1]*nshape[1], size[0]*nshape[2], size[0]*nshape[3]]
                    slide_layout = layout[i]['layout_results'][0]
                    center = ((shape[0]+shape[1])/2, (shape[2]+shape[3])/2)  # (x,y)
                    text_type = "Unknown"
                    for j in range(len(slide_layout)):
                        coordinates = slide_layout[j]['coordinate']  # [x1,y1,x2,y2]
                        if center[0] >= coordinates[0] \
                            and center[0] <= coordinates[2] \
                            and center[1] >= coordinates[1] \
                                and center[1] <= coordinates[3]:
                            # Center is inside layout
                            text_type = slide_layout[j]['label']
                            text_type = paddle_text_types.get(text_type)
                            if text_type is None:
                                text_type = "Unknown"
                            break
                    if text_type == "Unknown":
                        # Try seeing if there is any overlap
                        def rectangles_overlap(xa1, xa2, ya1, ya2, xb1, xb2, yb1, yb2):
                            # Make sure the coordinates are ordered correctly
                            xa1, xa2 = min(xa1, xa2), max(xa1, xa2)
                            ya1, ya2 = min(ya1, ya2), max(ya1, ya2)
                            xb1, xb2 = min(xb1, xb2), max(xb1, xb2)
                            yb1, yb2 = min(yb1, yb2), max(yb1, yb2)

                            # Check for non-overlap
                            if xa2 <= xb1 or xb2 <= xa1:
                                return False  # No overlap in x-axis
                            if ya2 <= yb1 or yb2 <= ya1:
                                return False  # No overlap in y-axis

                            return True  # Overlap exists

                        for j in range(len(slide_layout)):
                            coordinates = slide_layout[j]['coordinate']  # [x1,y1,x2,y2]
                            if rectangles_overlap(shape[0], shape[1], shape[2], shape[3],
                                                coordinates[0], coordinates[2], coordinates[1], coordinates[3]):
                                # Center is inside layout
                                text_type = slide_layout[j]['label']
                                text_type = paddle_text_types.get(text_type)
                                if text_type is None:
                                    text_type = "Body"
                                break

                    slide_source.append({
                        "type": text_type,
                        "text": nshape[4].replace('\'', '’')
                    })
                source.append(slide_source)
            return source


        def extract_text_shapes(pptx_path: str):
            try:
                if not os.path.isfile(pptx_path):
                    return f"No english file found at: {pptx_path}"

                eng_pptx = Presentation(pptx_path)
                pptx_rects = [eng_pptx.slide_width / eng_pptx.slide_height]
                pptx_zip = zipfile.ZipFile(pptx_path, "r")

                AVERAGE_CHAR_WIDTH_FACTOR = 0.47
                LINE_SPACING_FACTOR = 1.0

                def emu_to_pt(emu):
                    return emu / 12700

                def pt_to_emu(pt):
                    return int(pt * 12700)

                def process_table(table_shape: GraphicFrame, group_left=0, group_top=0):
                    """Processes a table and extracts bounding boxes for each cell."""
                    table = table_shape.table
                    num_rows = len(table.rows)
                    num_cols = len(table.columns)

                    shape_left = table_shape.left + group_left
                    shape_top = table_shape.top + group_top
                    shape_width = table_shape.width
                    shape_height = table_shape.height

                    row_heights = [row.height for row in table.rows]
                    col_widths = [col.width for col in table.columns]

                    total_width = sum(col_widths)
                    total_height = sum(row_heights)

                    scale_factor = (
                        shape_height / total_height if total_height > 0 else 1
                    )

                    y_cursor = shape_top
                    for row_idx, row_height in enumerate(row_heights):
                        if row_height == 0:
                            for col_idx, _ in enumerate(col_widths):
                                cell = table.cell(row_idx, col_idx)
                                num_lines = 0
                                font_size = 0
                                for para in cell.text_frame.paragraphs:
                                    num_lines += 1
                                    for run in para.runs:
                                        if run.text.strip() and run.font.size:
                                            line_spacing_factor = para.line_spacing if para.line_spacing else LINE_SPACING_FACTOR
                                            font_size = max(font_size, run.font.size.pt * line_spacing_factor)
                            row_height = max(row_height, num_lines * pt_to_emu(font_size))
                        row_height = row_height * scale_factor
                        x_cursor = shape_left
                        for col_idx, col_width in enumerate(col_widths):
                            cell = table.cell(row_idx, col_idx)

                            if not cell.text.strip():
                                x_cursor += col_width
                                continue  # Skip empty cells

                            x1 = x_cursor / eng_pptx.slide_width
                            x2 = (x_cursor + col_width) / eng_pptx.slide_width
                            y1 = y_cursor / eng_pptx.slide_height
                            y2 = (y_cursor + row_height) / eng_pptx.slide_height

                            yield [x1, x2, y1, y2, cell.text.strip()]
                            x_cursor += col_width

                        y_cursor += row_height

                def process_shape(shape, x_scale=1, y_scale=1, x_transform=0, y_transform=0, initial_left=-1, initial_top=-1):
                    """ Processes a shape or a grouped shape recursively. """
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        shape: GroupShape = shape

                        shape_left = shape.left + x_transform
                        shape_top = shape.top + y_transform
                        if initial_left != -1 and initial_top != -1:
                            shape_left = shape_left + (shape.left - initial_left) * (x_scale - 1)
                            shape_top = shape_top  + (shape.top - initial_top) * (y_scale - 1)

                        minX1 = math.inf
                        maxX2 = 0
                        minY1 = math.inf
                        maxY2 = 0
                        for sub_shape in shape.shapes:
                            if sub_shape.left < minX1: minX1 = sub_shape.left
                            if sub_shape.left + sub_shape.width > maxX2: maxX2 = sub_shape.left + sub_shape.width
                            if sub_shape.top < minY1: minY1 = sub_shape.top
                            if sub_shape.top + sub_shape.height > maxY2: maxY2 = sub_shape.top + sub_shape.height
                        initial_group_width = maxX2 - minX1
                        initial_group_height = maxY2 - minY1
                        if initial_group_width == 0: group_scale_x = 1
                        else: group_scale_x = shape.width / initial_group_width
                        if initial_group_height == 0: group_scale_y = 1
                        else: group_scale_y = shape.height / initial_group_height
                        group_transform_x = shape_left - minX1
                        group_transform_y = shape_top - minY1
                        for sub_shape in shape.shapes:
                            yield from process_shape(sub_shape, group_scale_x, group_scale_y, group_transform_x, group_transform_y, minX1, minY1)
                            # yield from process_shape(sub_shape)
                        return

                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        # yield from process_table(shape, group_left, group_top)
                        yield from process_table(shape)
                        return

                    elif not shape.has_text_frame:
                        return

                    text_frame = shape.text_frame
                    if not text_frame.paragraphs:
                        return

                    # Adjust shape dimensions based on group offsets
                    shape_width_pt = emu_to_pt(shape.width)
                    shape_height_pt = emu_to_pt(shape.height)

                    para_info = []
                    total_estimated_height = 0
                    first = True

                    for para in text_frame.paragraphs:
                        text = para.text.strip()
                        is_empty = not text

                        font_size = 0
                        for run in para.runs:
                            if run.text.strip() and run.font.size:
                                font_size = max(font_size, run.font.size.pt)
                        if font_size == 0:
                            font_size = 12

                        line_spacing_factor = para.line_spacing if para.line_spacing else LINE_SPACING_FACTOR

                        if is_empty:
                            num_lines = 1
                        else:
                            char_width_pt = font_size * AVERAGE_CHAR_WIDTH_FACTOR
                            chars_per_line = max(1, math.floor(shape_width_pt / char_width_pt))
                            num_lines = max(1, math.ceil(len(text) / chars_per_line))

                        para_height_pt = (
                            1.3 * num_lines * font_size * line_spacing_factor 
                            - 0.25 * (num_lines - 1) * font_size * line_spacing_factor 
                            + emu_to_pt(para.space_before if para.space_before else 0) 
                            + emu_to_pt(para.space_after if para.space_after else 0)
                        )
                        para_info.append({'para': para, 'height_pt': para_height_pt, 'is_empty': is_empty})
                        total_estimated_height += para_height_pt
                        first = False

                    scale_factor = (
                        (shape_height_pt - emu_to_pt(text_frame.margin_bottom + text_frame.margin_top)) / total_estimated_height
                        if total_estimated_height > 0 else 1
                    )

                    shape_left = shape.left + x_transform
                    shape_top = shape.top + text_frame.margin_top + y_transform
                    if initial_left != -1 and initial_top != -1:
                        shape_left = shape_left + (shape.left - initial_left) * (x_scale - 1)
                        shape_top = shape_top  + (shape.top - initial_top) * (y_scale - 1)
                    y_cumulative = shape_top

                    for info in para_info:
                        scaled_height_pt = info['height_pt'] * scale_factor

                        if info['is_empty']:
                            y_cumulative += pt_to_emu(scaled_height_pt)
                            continue

                        scaled_height_emu = pt_to_emu(scaled_height_pt)

                        x1 = shape_left / eng_pptx.slide_width
                        x2 = (shape_left + shape.width * x_scale) / eng_pptx.slide_width
                        y1 = y_cumulative / eng_pptx.slide_height
                        y2 = (y_cumulative + scaled_height_emu * y_scale) / eng_pptx.slide_height

                        yield [x1, x2, y1, y2, info['para'].text.strip()]
                        y_cumulative += scaled_height_emu

                def is_slide_hidden(slide_number):
                    """Checks if a specific slide is hidden based on the 'show' attribute in its XML."""
                    slide_path = f"ppt/slides/slide{slide_number}.xml"

                    with pptx_zip.open(slide_path) as xml_file:
                        tree = etree.parse(xml_file)

                    # Get the root element <p:sld>
                    root = tree.getroot()

                    # Check if 'show' attribute is set to "0"
                    return root.get("show") == "0"

                for i, slide in enumerate(eng_pptx.slides):
                    if is_slide_hidden(i + 1): continue
                    slide_paragraph_rects = []
                    for shape in slide.shapes:
                        slide_paragraph_rects.extend(process_shape(shape))
                    pptx_rects.append(slide_paragraph_rects)

                return pptx_rects

            except Exception as e:
                traceback.print_exc()
                return str(e)
        

        logger.publish("Analyzing text shapes in pptx file...")
        extracted_shapes = extract_text_shapes(ppt_path)
        scale_factor = extracted_shapes[0]
        extracted_shapes.pop(0)
        number_of_slides = len(extracted_shapes)

        logger.publish("Processing layout types...")
        layout = analyze_document_layout(pdf_path, number_of_slides)
        os.remove(pdf_path)

        logger.publish("Applying layout types to text shapes...")
        source = apply_layout_types(extracted_shapes, layout)
        logger.publish("Analyzing shapes complete.")

        return source, number_of_slides
